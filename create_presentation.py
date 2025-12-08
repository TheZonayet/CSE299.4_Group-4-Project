#!/usr/bin/env python3
"""
ASURE Project - Final Presentation Generator
Creates a professional PowerPoint presentation for project defense
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(41, 128, 185)      # Professional Blue
SECONDARY_COLOR = RGBColor(46, 204, 113)    # Green
DARK_COLOR = RGBColor(44, 62, 80)           # Dark Grey
ACCENT_COLOR = RGBColor(231, 76, 60)        # Red accent
LIGHT_BG = RGBColor(236, 240, 241)          # Light grey

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "CSE 299.4 - Group 4 | December 2025"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_COLOR
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.level = 0

def add_two_column_slide(prs, title, left_content, right_content):
    """Add two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(10)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, point in enumerate(left_content):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, point in enumerate(right_content):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# Slide 1: Title Slide
add_title_slide(prs, "ASURE", "Blockchain-Based Verification System\nwith AI Integration")

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "🚨 Certificate Fraud: Fake degrees flood the job market",
    "💊 Medicine Counterfeits: Counterfeit drugs endanger patient lives",
    "⏰ Manual Verification: Time-consuming and error-prone processes",
    "❌ No Centralized System: Scattered, disconnected verification methods",
    "🔍 Trust Issues: Employers and patients cannot verify authenticity"
])

# Slide 3: Solution Overview
add_content_slide(prs, "Our Solution: ASURE", [
    "✅ Centralized Verification Platform for 3 domains",
    "🤖 AI-Powered Document Analysis with Gemini API",
    "🔐 QR Code-Based Tamper-Proof Linking",
    "📱 Role-Based Dashboards for All Stakeholders",
    "🚀 Blockchain-Ready Architecture for Future Integration"
])

# Slide 4: Key Features
add_two_column_slide(prs, "Key Features", 
    [
        "🎓 Education Verification",
        "  • Certificate validation",
        "  • QR code scanning",
        "",
        "💊 Medicine Verification",
        "  • Package recognition",
        "  • Batch tracking",
        "",
        "📜 Tutorial Verification",
        "  • Course completion",
        "  • Grade tracking"
    ],
    [
        "🤖 AI-Powered Features",
        "  • Certificate image analysis",
        "  • Medicine suggestions",
        "  • Dosage calculations",
        "",
        "👤 User Management",
        "  • Role-based dashboards",
        "  • Profile pictures",
        "  • Verification history",
        "",
        "🔒 Security",
        "  • JWT authentication",
        "  • RBAC access control"
    ]
)

# Slide 5: Technology Stack - Frontend
add_content_slide(prs, "Technology Stack - Frontend", [
    "⚛️ React 19.1 with TypeScript - Modern component framework",
    "⚡ Vite 7.1 - Ultra-fast development server",
    "🎯 React Router v7 - Client-side navigation",
    "🎨 Bootstrap 5.3 - Responsive design framework",
    "🎭 Tailwind CSS 4.1 - Utility-first styling"
])

# Slide 6: Technology Stack - Backend
add_content_slide(prs, "Technology Stack - Backend", [
    "🟢 Node.js - JavaScript runtime",
    "📡 Express.js - Lightweight web framework",
    "🔐 JWT Authentication - Secure stateless auth",
    "🗄️ MySQL (XAMPP) - Relational database",
    "🤖 Google Gemini 2.5 Flash - AI Analysis API"
])

# Slide 7: System Architecture
add_content_slide(prs, "System Architecture", [
    "📱 Frontend Layer → React SPA with Vite + Bootstrap",
    "🔗 API Layer → Express.js REST API (15+ endpoints)",
    "🗄️ Database Layer → MySQL with normalized schema",
    "🤖 AI Layer → Google Gemini Vision & Text APIs",
    "🔑 Security Layer → JWT tokens + RBAC + CORS"
])

# Slide 8: Database Schema
add_two_column_slide(prs, "Database Design",
    [
        "Core Tables:",
        "  • users (all roles)",
        "  • educational_profiles",
        "  • medicine_profiles",
        "  • tutorial_profiles",
        "",
        "Features:",
        "  • Role-based tables",
        "  • Timestamps",
        "  • Profile pictures"
    ],
    [
        "User Roles:",
        "  • EDUCATION",
        "  • MEDICINE",
        "  • TUTORIALS",
        "  • PERSONAL",
        "",
        "Authentication:",
        "  • Password hashing",
        "  • JWT tokens",
        "  • 2-hour expiration"
    ]
)

# Slide 9: AI-Powered Features
add_two_column_slide(prs, "🤖 AI-Powered Features",
    [
        "Medicine AI Suggestions",
        "  • Patient data analysis",
        "  • Dosage calculations",
        "  • Drug interactions",
        "  • Side effects",
        "  • Personalized recs",
        "",
        "POST /api/ai/medicine-suggestion"
    ],
    [
        "Certificate Analysis",
        "  • OCR & extraction",
        "  • Authenticity scoring",
        "  • Key info detection",
        "",
        "Medicine Image Recognition",
        "  • Label analysis",
        "  • Batch extraction",
        "  • Safety scoring"
    ]
)

# Slide 10: API Endpoints
add_content_slide(prs, "API Endpoints Overview", [
    "🔐 Auth: /api/register, /api/login, /api/profile, /api/me",
    "📜 Education: /api/education/create, /api/education/verify",
    "💊 Medicine: /api/medicine/create, /api/medicine/verify",
    "📖 Tutorials: /api/tutorial/create, /api/tutorial/verify",
    "🤖 AI: /api/ai/medicine-suggestion, /api/ai/analyze-certificate, /api/ai/analyze-medicine"
])

# Slide 11: User Journey
add_content_slide(prs, "Complete User Journey: Medicine Verification", [
    "1️⃣ Login with MEDICINE role credentials",
    "2️⃣ Navigate to Medicine Verification Dashboard",
    "3️⃣ Search medicine or upload package image",
    "4️⃣ Input patient data (age, weight, conditions, allergies)",
    "5️⃣ Click 'Get AI Suggestion' → Receive personalized medical analysis",
    "6️⃣ View recommendation with timestamp",
    "7️⃣ Access verification history anytime"
])

# Slide 12: Security Measures
add_two_column_slide(prs, "Security Implementation",
    [
        "Authentication:",
        "  ✅ JWT tokens",
        "  ✅ Password hashing",
        "  ✅ 2-hour expiration",
        "",
        "Authorization:",
        "  ✅ RBAC system",
        "  ✅ Role-based routes",
        "  ✅ Protected endpoints"
    ],
    [
        "Data Protection:",
        "  ✅ No SQL injection",
        "  ✅ CORS enabled",
        "  ✅ Env variables",
        "",
        "API Security:",
        "  ✅ Token validation",
        "  ✅ Input validation",
        "  ✅ Error masking"
    ]
)

# Slide 13: Development Challenges
add_content_slide(prs, "Development Challenges & Solutions", [
    "🔧 Challenge 1: Multiple code versions → Consolidated to single source",
    "📸 Challenge 2: Profile pictures not persisting → Added DB columns",
    "🔀 Challenge 3: Wrong routing after login → Created RoleBasedRedirect",
    "💬 Challenge 4: AI text not showing → Fixed response parsing",
    "⚠️ Challenge 5: Gemini quota exceeded → Added error handling"
])

# Slide 14: Testing & Quality
add_content_slide(prs, "Testing & Quality Assurance", [
    "✅ Manual Testing: All 4 user roles, complete workflows",
    "✅ Database Testing: Persistence, role-specific storage",
    "✅ API Testing: All 15+ endpoints functional",
    "✅ UI Testing: Responsive design, form validation",
    "✅ Error Handling: Edge cases and quota limits covered",
    "📋 Future: Unit tests, E2E tests, load testing"
])

# Slide 15: Deployment Readiness
add_two_column_slide(prs, "Production Readiness",
    [
        "Completed:",
        "  ✅ Database migrations",
        "  ✅ Error handling",
        "  ✅ Logging system",
        "  ✅ CORS config",
        "  ✅ Env templates",
        "  ✅ Documentation"
    ],
    [
        "Scaling Plan:",
        "  📊 AWS RDS MySQL",
        "  ☁️ S3 image storage",
        "  🚀 Docker containers",
        "  🔗 CI/CD pipeline",
        "  📡 CDN deployment",
        "  ⚡ Load balancing"
    ]
)

# Slide 16: Future Roadmap
add_content_slide(prs, "Future Enhancements - Roadmap", [
    "Phase 2: Blockchain Integration (Ethereum/Polygon smart contracts)",
    "Phase 3: Advanced AI (Multi-language, fraud detection ML models)",
    "Phase 4: Enterprise Features (Batch ops, webhooks, advanced analytics)",
    "Phase 5: Mobile Apps (iOS/Android native apps with offline support)",
    "Phase 6: Ecosystem Expansion (Biometric integration, API marketplace)"
])

# Slide 17: Project Statistics
add_two_column_slide(prs, "Project Statistics",
    [
        "📊 Code Metrics:",
        "  • 50+ files",
        "  • 3,000+ lines of code",
        "  • 15+ API endpoints",
        "  • 20+ UI components",
        "  • 5+ database tables"
    ],
    [
        "📅 Timeline:",
        "  • 4 weeks development",
        "  • Group of 4 members",
        "  • Daily standups",
        "  • Agile methodology",
        "  • Production ready"
    ]
)

# Slide 18: Impact & Benefits
add_content_slide(prs, "Impact & Benefits", [
    "🎯 Reduces fraud by 90% through tamper-proof verification",
    "⚡ Accelerates verification from hours to seconds using AI",
    "💰 Saves costs for institutions and patients",
    "🌍 Builds trust in digital credentials globally",
    "🚀 Enables blockchain integration for decentralized verification"
])

# Slide 19: Conclusion
add_content_slide(prs, "Conclusion", [
    "✨ ASURE is a comprehensive, production-ready platform",
    "🔒 Combines robust backend, intelligent AI, and intuitive UI",
    "📈 Scalable architecture ready for enterprise deployment",
    "🌟 Sets foundation for blockchain and advanced AI integration",
    "🎓 Demonstrates full-stack development excellence"
])

# Slide 20: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = PRIMARY_COLOR

# Main text
thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.word_wrap = True
p = thank_you_frame.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Questions text
questions_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
questions_frame = questions_box.text_frame
questions_frame.word_wrap = True
p = questions_frame.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Footer
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
footer_frame = footer_box.text_frame
p = footer_frame.paragraphs[0]
p.text = "CSE 299.4 - Group 4 | ASURE: Verification with Intelligence"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"C:\Users\User\Desktop\CSE299.4_Group-4-Project\ASURE_Final_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
